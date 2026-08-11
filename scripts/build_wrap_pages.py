#!/usr/bin/env python3
"""build_wrap_pages.py - 为任意 PPT 添加豪华版封面/目录/结束页

用法（模式 C 增强 / build_pptx.py 之后）:
    python3 scripts/build_wrap_pages.py <input.pptx> --title "主标题" \\
        [--subtitle "副标题"] [--tagline "英文衬底"] \\
        [--dept 部门] [--version 版本] [--date 日期] \\
        [--toc '{"01":"章节一","02":"章节二"}' ] [--toc-sub '{"01":"副标一"}' ] \\
        [--accent deepblue|red|gold|green] [--out out.pptx]

说明:
- 用 deepcopy 方案把 input 的内容页复制到全新 PPT，避免 build_pptx.py 残留
  parts 导致的 part 命名冲突（soffice 加载失败）。
- 追加 3 张豪华页（渐变背景 + 几何装饰 + 图标系统 + 强对比标题），
  按「封面 -> 目录 -> 内容 -> 结束」重排。
- 主题色 --accent:
    deepblue  #1F3A93/#0E1B42  (architecture-deck 系)
    red       #A52524/#5E1212  (premium-corp 系)
    gold      #8A6D1F/#3E320E  (政务/高端商务)
    green     #1E5B4C/#0F2E26  (环保/科技)
"""
import argparse
import copy
import json
import sys

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ACCENTS = {
    "deepblue": {"primary": (0x1F, 0x3A, 0x93), "dark": (0x0E, 0x1B, 0x42), "mid": (0x2A, 0x4F, 0xB8), "gold": (0xC9, 0x9D, 0x52)},
    "red":     {"primary": (0xA5, 0x25, 0x24), "dark": (0x5E, 0x12, 0x12), "mid": (0xC9, 0x3A, 0x37), "gold": (0xE8, 0xC5, 0x7E)},
    "gold":    {"primary": (0x8A, 0x6D, 0x1F), "dark": (0x3E, 0x32, 0x0E), "mid": (0xB8, 0x95, 0x35), "gold": (0xE8, 0xD9, 0x9E)},
    "green":   {"primary": (0x1E, 0x5B, 0x4C), "dark": (0x0F, 0x2E, 0x26), "mid": (0x2E, 0x7D, 0x69), "gold": (0xC9, 0xDD, 0x9A)},
}


def rgb(t):
    return RGBColor(*t)


def add_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def add_oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(l), Cm(t), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = "Microsoft YaHei"; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_line(slide, l1, t1, l2, t2, color, width=2):
    ln = slide.shapes.add_connector(1, Cm(l1), Cm(t1), Cm(l2), Cm(t2))
    ln.line.color.rgb = color; ln.line.width = Pt(width)
    return ln


def gradient_bg(slide, c1, c2, width, height, angle=90):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    fill = bg.fill; fill.gradient(); fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].color.rgb = c1; stops[0].position = 0.0
    stops[1].color.rgb = c2; stops[1].position = 1.0
    bg.line.fill.background()
    return bg


def build_wrap(input_pptx, out_pptx, title, subtitle, tagline,
               dept, version, date, toc, toc_sub, accent, highlights):
    """核心流程：deepcopy 内容页 + 3 张豪华辅助页 + 重排"""
    a = ACCENTS[accent]
    P, D, M, G = rgb(a["primary"]), rgb(a["dark"]), rgb(a["mid"]), rgb(a["gold"])
    W, L, GR = RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xF5, 0xF1, 0xEC), RGBColor(0x55, 0x55, 0x55)
    DG = RGBColor(0x2B, 0x2B, 0x2B)

    src = Presentation(input_pptx)
    n = len(src.slides)

    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # 1) 复制内容页
    for slide in src.slides:
        ns = prs.slides.add_slide(blank)
        for shape in slide.shapes:
            ns.shapes._spTree.append(copy.deepcopy(shape._element))

    # 2) 封面
    cover = prs.slides.add_slide(blank)
    gradient_bg(cover, D, P, prs.slide_width, prs.slide_height)
    add_oval(cover, 25, 1, 8, 8, M)
    add_oval(cover, 27, 3, 4, 4, P)
    add_oval(cover, -1, 14, 5, 5, M)
    add_oval(cover, 31, 13, 1.2, 1.2, G)
    add_oval(cover, 29, 16, 0.6, 0.6, G)
    add_rect(cover, 2, 1.0, 8, 0.18, G)
    add_rect(cover, 2, 3.0, 0.12, 7.0, G)
    if tagline:
        add_text(cover, 3.0, 1.0, 18, 0.6, tagline, size=13, bold=True, color=G)
    add_text(cover, 3, 3.6, 18, 1.8, title, size=42, bold=True, color=W)
    if subtitle:
        add_text(cover, 3, 5.6, 18, 1.2, subtitle, size=26, color=RGBColor(0xF5, 0xD9, 0xD8))
    add_line(cover, 3, 8.0, 18, 8.0, G, width=3)
    info = [("▣ 部门", dept), ("◉ 版本", version), ("◈ 日期", date)]
    for i, (icon, val) in enumerate(info):
        x = 3 + i * 6.0
        add_rect(cover, x, 9.2, 5.6, 1.5, D)
        add_text(cover, x + 0.5, 9.4, 5, 0.5, icon, size=11, bold=True, color=G)
        add_text(cover, x + 0.5, 9.95, 5, 0.6, val, size=12, bold=True, color=W)
    if highlights:
        add_text(cover, 19, 9.5, 12, 0.6, "HIGHLIGHTS", size=10, bold=True, color=G)
        add_line(cover, 19, 10.1, 31, 10.1, G, width=1)
        for i, t in enumerate(highlights[:4]):
            y = 10.4 + i * 0.75
            add_text(cover, 19, y, 1, 0.6, "◆", size=10, color=G)
            add_text(cover, 20.2, y, 12, 0.6, t, size=12, color=W)
    add_rect(cover, 27.5, 1.0, 4.5, 1.2, D)
    add_text(cover, 27.5, 1.25, 4.5, 0.7, "⚙  LOGO", size=14, bold=True, color=G, align=PP_ALIGN.CENTER)
    add_rect(cover, 0, 18.3, 33.87, 0.75, D)
    add_text(cover, 2, 18.45, 30, 0.5, "AGENT POWER · 智能编制  |  仅供学习与研究使用", size=9, color=RGBColor(0xD9, 0xB4, 0xB3))

    # 3) 目录
    agenda = prs.slides.add_slide(blank)
    gradient_bg(agenda, W, L, prs.slide_width, prs.slide_height)
    add_rect(agenda, 0, 0, 33.87, 0.5, P)
    add_rect(agenda, 0, 0.5, 33.87, 0.06, G)
    add_rect(agenda, 2, 1.2, 0.7, 0.4, G)
    add_text(agenda, 2.9, 1.1, 8, 0.5, "CONTENTS", size=12, bold=True, color=P)
    add_text(agenda, 2, 1.9, 8, 1.0, "目录", size=34, bold=True, color=DG)
    add_text(agenda, 2, 3.0, 8, 0.5, f"{len(toc)} SECTIONS", size=11, bold=True, color=G)
    add_line(agenda, 2, 3.6, 8, 3.6, G, width=1)
    for i, (num, t) in enumerate(toc.items()):
        col, row = i % 2, i // 2
        x = 10 + col * 12
        y = 1.6 + row * 2.6
        add_rect(agenda, x, y, 1.5, 1.5, P)
        add_text(agenda, x, y + 0.35, 1.5, 0.9, num, size=18, bold=True, color=W, align=PP_ALIGN.CENTER)
        add_text(agenda, x + 1.8, y + 0.1, 9.5, 0.8, t, size=16, bold=True, color=DG)
        sub = toc_sub.get(num, "") if toc_sub else ""
        if sub:
            add_text(agenda, x + 1.8, y + 0.95, 9.5, 0.5, sub, size=10, color=GR)
        add_line(agenda, x + 1.8, y + 1.45, x + 11, y + 1.45, color=RGBColor(0xE8, 0xD8, 0xD5), width=1)
        add_oval(agenda, x + 11.3, y + 0.4, 0.5, 0.5, G)
    add_rect(agenda, 22, 17.4, 10, 0.05, P)
    add_text(agenda, 22, 17.6, 10, 0.4, "AGENT POWER · 2026", size=9, bold=True, color=P)
    add_text(agenda, 22, 18.0, 10, 0.4, "PROJECT INTRODUCTION", size=9, color=GR)

    # 4) 结束页
    ending = prs.slides.add_slide(blank)
    gradient_bg(ending, P, D, prs.slide_width, prs.slide_height)
    for cx, cy, L, w in [(1.2, 1.2, 2.0, 0.08), (32.5, 1.2, 2.0, 0.08), (1.2, 17.6, 2.0, 0.08), (32.5, 17.6, 2.0, 0.08)]:
        add_rect(ending, cx, cy, L, w, G)
        add_rect(ending, cx, cy, w, L, G)
    add_oval(ending, 4.0, 1.5, 1.0, 1.0, G)
    add_oval(ending, 28.5, 15.5, 1.2, 1.2, M)
    add_oval(ending, 30.2, 17.0, 0.6, 0.6, G)
    add_text(ending, 0, 14.5, 7, 3, "END", size=80, bold=True, color=D, align=PP_ALIGN.CENTER)
    add_text(ending, 8, 4.0, 22, 1.8, "感谢聆听", size=56, bold=True, color=W, align=PP_ALIGN.CENTER)
    add_text(ending, 8, 6.0, 22, 0.9, "THANKS FOR LISTENING", size=14, color=RGBColor(0xE8, 0xC5, 0xC4), align=PP_ALIGN.CENTER)
    add_line(ending, 13, 7.4, 21, 7.4, G, width=3)
    if subtitle:
        add_rect(ending, 4, 8.0, 25.8, 2.8, D)
        add_text(ending, 5, 8.3, 24, 0.6, "▸ 下一步计划 NEXT STEPS", size=12, bold=True, color=G)
        add_text(ending, 5, 8.9, 24, 1.6, "敬请指导 · 适时修订 · 持续优化落地", size=13, color=W)
    add_rect(ending, 4, 11.4, 25.8, 0.05, G)
    contacts = [("✉ 邮箱", "support@agent-ppt.com"), ("✆ 微信", "Yu_Z0809"), ("⌂ 主页", "github.com/lmori1301"), ("◈ 版本", version)]
    for i, (l, v) in enumerate(contacts):
        x = 4.5 + i * 6.5
        add_text(ending, x, 11.8, 6, 0.5, l, size=10, color=RGBColor(0xE8, 0xC5, 0xC4), align=PP_ALIGN.CENTER)
        add_text(ending, x, 12.3, 6, 0.5, v, size=12, bold=True, color=W, align=PP_ALIGN.CENTER)
    add_rect(ending, 4, 13.1, 25.8, 0.05, G)
    add_rect(ending, 0, 17.6, 33.87, 1.45, D)
    add_text(ending, 2, 17.9, 30, 0.4, "PROJECT · AGENT POWER", size=10, bold=True, color=G)
    add_text(ending, 2, 18.5, 30, 0.4, f"{title} · 智能编制 · 仅供学习与研究使用", size=9, color=RGBColor(0xD9, 0xB4, 0xB3))

    # 5) 重排：封面, 目录, 内容(0..n-1), 结束
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    content = slides[:n]
    cover_el, agenda_el, ending_el = slides[n], slides[n + 1], slides[n + 2]
    for el in list(xml_slides):
        xml_slides.remove(el)
    for el in [cover_el, agenda_el] + content + [ending_el]:
        xml_slides.append(el)

    prs.save(out_pptx)
    print(f"已生成豪华版包装页：{out_pptx}（{n} 内容页 + 封面/目录/结束 = {n + 3} 页）")


def main():
    ap = argparse.ArgumentParser(description="为已有 PPT 添加豪华版封面/目录/结束页")
    ap.add_argument("input", help="输入 PPTX 路径")
    ap.add_argument("--title", required=True, help="主标题")
    ap.add_argument("--subtitle", default="", help="副标题")
    ap.add_argument("--tagline", default="", help="英文衬底标签")
    ap.add_argument("--dept", default="项目组", help="部门")
    ap.add_argument("--version", default="V1.0", help="版本")
    ap.add_argument("--date", default="2026 年", help="日期")
    ap.add_argument("--toc", default="{}", help="JSON 目录 {01:标题, 02:标题}")
    ap.add_argument("--toc-sub", default="{}", help="JSON 目录副标 {01:副标}")
    ap.add_argument("--highlights", default="", help="封面亮点列表，分号分隔（最多 4 项）")
    ap.add_argument("--accent", default="deepblue", choices=list(ACCENTS.keys()), help="主题色")
    ap.add_argument("--out", required=True, help="输出 PPTX 路径")
    args = ap.parse_args()

    toc = json.loads(args.toc)
    toc_sub = json.loads(args.toc_sub)
    highlights = [h.strip() for h in args.highlights.split(";") if h.strip()] if args.highlights else []
    build_wrap(args.input, args.out, args.title, args.subtitle, args.tagline,
               args.dept, args.version, args.date, toc, toc_sub, args.accent, highlights)


if __name__ == "__main__":
    main()